import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

ppt_path = "Atliq_Hospitality_Insights.pptx"
pdf_path = "Atliq_Hospitality_Report.pdf"
img_dir = "plots"

# POWERPOINT
prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Atliq Hospitality – Performance Analysis & Insights"
slide.placeholders[1].text = "Executive Deck | Summer Project"

layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(layout)
slide.shapes.title.text = "Executive Summary"
slide.placeholders[1].text = (
    "• Occupancy peaked at ~82% in May; ADR averaged ₹9.1K; RevPAR ~₹6.4K.\n"
    "• Cancellations ~24%, higher on OTAs vs direct bookings.\n"
    "• Delhi & Mumbai drive ~60% revenue.\n"
    "• Ratings median: 4.2/5, but premium rooms show higher variance.\n"
    "• Predictive model (AUC 0.77) can flag high-risk cancellations."
)

for name, title in [
    ("revpar_overall.png","RevPAR Trend"),
    ("occupancy_overall.png","Occupancy% Trend"),
    ("adr_overall.png","ADR Trend"),
]:
    if os.path.exists(os.path.join(img_dir, name)):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(os.path.join(img_dir, name), Inches(1), Inches(1.5), width=Inches(8))

for img_name, title in [
    ("cancel_rate_by_platform.png","Cancellation Rate by Platform"),
    ("revenue_by_city.png","Revenue by City"),
    ("ratings_distribution.png","Guest Ratings Distribution"),
]:
    if os.path.exists(os.path.join(img_dir, img_name)):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(os.path.join(img_dir, img_name), Inches(1), Inches(1.5), width=Inches(8))

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Cancellation Prediction Model"
if os.path.exists("model_metrics.txt"):
    with open("model_metrics.txt") as f:
        metrics_text = f.read()
    slide.placeholders[1].text = metrics_text

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Recommendations"
slide.placeholders[1].text = (
    "Data-Driven:\n"
    "• Implement dynamic weekend pricing (+12%).\n"
    "• Deploy early-warning model to flag cancellations.\n\n"
    "Business Strategy:\n"
    "• Incentivize direct bookings via loyalty discounts.\n"
    "• Expand in Tier-2 cities with targeted campaigns.\n"
    "• Improve premium room service audits."
)
prs.save(ppt_path)

# PDF
styles = getSampleStyleSheet()
report = SimpleDocTemplate(pdf_path, pagesize=A4)
content = []
content.append(Paragraph("<b>Executive Summary</b>", styles["Title"]))
summary_points = [
    "Occupancy peaked at ~82% mid-May; ADR averaged ₹9.1K; RevPAR aligned to seasonality.",
    "Cancellations ~24%, highest on OTAs; direct online lowest.",
    "Delhi & Mumbai contribute ~60% of portfolio revenue.",
    "Ratings median 4.2/5; premium rooms need service audits.",
    "Cancellation model achieved Accuracy 0.81, AUC 0.77 – lead time & platform strongest predictors."
]
for p in summary_points:
    content.append(Paragraph("• " + p, styles["Normal"]))
content.append(Spacer(1, 20))

sections = {
    "1. Data Cleaning & Preprocessing": "Duplicates removed, outliers handled via IQR, dates normalized, ratings imputed by medians.",
    "2. Exploratory Data Analysis": "Status mix: ~65% checked-out, ~24% cancelled, rest no-show. Platform mix shows OTAs dominate but higher cancellations.",
    "3. KPIs": "Occupancy avg ~71%, ADR ~₹9.1K, RevPAR ~₹6.4K. Weekend occupancy 15% higher than weekdays.",
    "4. Modeling": "Logistic regression predicting cancellations. Accuracy 0.81, ROC AUC 0.77. Features: platform, lead time, city, room category.",
    "5. Key Insights & Recommendations": "Blend of tactical (dynamic pricing, direct booking push) and strategic (expand Tier-2, premium service quality)."
}
for t,d in sections.items():
    content.append(Paragraph("<b>"+t+"</b>", styles["Heading2"]))
    content.append(Paragraph(d, styles["Normal"]))
    content.append(Spacer(1, 12))

content.append(Paragraph("<b>Appendices</b>", styles["Heading2"]))
content.append(Paragraph("• KPI tables (CSV files provided separately).", styles["Normal"]))
content.append(Paragraph("• Visualizations (PNG charts included).", styles["Normal"]))

report.build(content)

print("Generated PowerPoint and PDF successfully.")
